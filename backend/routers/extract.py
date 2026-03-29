"""File extraction endpoint — converts uploaded documents/images to policy text."""

from __future__ import annotations

import base64
import io
import logging

from fastapi import APIRouter, HTTPException, UploadFile
from langchain_core.messages import HumanMessage

from graph.llm import get_llm

logger = logging.getLogger(__name__)

router = APIRouter()

_VISION_PROMPT = (
    "This image was uploaded as part of an economic policy document. "
    "Extract and transcribe all visible text verbatim. If the image contains charts, "
    "tables, or diagrams, describe their key data and what they communicate. "
    "Be thorough — every detail may be relevant for policy analysis."
)


async def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(p.strip() for p in pages if p.strip())


async def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = " ".join(
                    para.text for para in shape.text_frame.paragraphs
                ).strip()
                if text:
                    parts.append(text)
    return "\n\n".join(parts)


async def _extract_image(data: bytes, mime: str) -> str:
    llm = get_llm(max_tokens=2048)
    b64 = base64.b64encode(data).decode()
    msg = HumanMessage(content=[
        {"type": "text", "text": _VISION_PROMPT},
        {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
    ])
    response = await llm.ainvoke([msg])
    return str(response.content)


@router.post("/extract")
async def extract_file(file: UploadFile) -> dict[str, str]:
    data = await file.read()
    filename = (file.filename or "").lower()
    content_type = file.content_type or ""

    logger.info("extract: file=%s  size=%d  type=%s", filename, len(data), content_type)

    try:
        if filename.endswith((".txt", ".md")):
            text = data.decode("utf-8", errors="replace")
        elif filename.endswith(".pdf") or "pdf" in content_type:
            text = await _extract_pdf(data)
        elif filename.endswith((".pptx", ".ppt")) or "presentation" in content_type:
            text = await _extract_pptx(data)
        elif filename.endswith((".png", ".jpg", ".jpeg", ".gif", ".webp")) or content_type.startswith("image/"):
            text = await _extract_image(data, content_type or "image/png")
        else:
            raise HTTPException(status_code=415, detail=f"Unsupported file type: {filename}")
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("extract: failed for %s", filename)
        raise HTTPException(status_code=500, detail=f"Extraction failed: {e}") from e

    if not text.strip():
        raise HTTPException(status_code=422, detail="No text could be extracted from the file.")

    logger.info("extract: extracted %d chars from %s", len(text), filename)
    return {"text": text.strip()}
